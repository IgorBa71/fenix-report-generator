# -*- coding: utf-8 -*-
"""
pptx_slides.py — отдельные слайды Презентации (python-pptx). Каждая функция
build_slide_<key>(pres, data) добавляет один слайд в pres и ничего не
возвращает. data — тот же JSON/dict, что выдаёт pptx_data_export.py.
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from pptx_presentation_builder import (
    add_slide, add_textbox, add_bullets, add_rounded_rect, add_line_arrow,
    title, footer, NAVY, TERRACOTTA, GOLD, TAUPE, TEAL, WHITE, OFFWHITE, DARKTEXT,
    FONT, W, H, BASE,
)


def build_slide_cover(pres, data):
    client = data["client"]
    slide = add_slide(pres, NAVY)
    logo_w = Inches(3.0)
    logo_h = int(logo_w * (221 / 357))
    slide.shapes.add_picture(
        str(__import__("pptx_presentation_builder").LOGO_VERTICAL),
        (W - logo_w) // 2, Inches(0.15), width=logo_w, height=logo_h,
    )
    add_textbox(slide, "РЕЗУЛЬТАТЫ ДИАГНОСТИКИ", Inches(1.5), Inches(2.2), W - Inches(3), Inches(0.7),
                size=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "И ПЛАН ДЕЙСТВИЙ", Inches(1.5), Inches(2.85), W - Inches(3), Inches(0.7),
                size=34, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

    line = slide.shapes.add_connector(1, W // 2 - Inches(0.9), Inches(3.85), W // 2 + Inches(0.9), Inches(3.85))
    line.line.color.rgb = TERRACOTTA
    line.line.width = Pt(2.5)

    add_textbox(slide, client["company"], Inches(1.5), Inches(4.15), W - Inches(3), Inches(0.5),
                size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, client["name"], Inches(1.5), Inches(4.65), W - Inches(3), Inches(0.4),
                size=15, color="B9C4CA", align=PP_ALIGN.CENTER)
    add_textbox(slide, client["date"], Inches(1.5), Inches(5.02), W - Inches(3), Inches(0.35),
                size=12, color="7C8B93", align=PP_ALIGN.CENTER)
    add_textbox(slide, "© ЛАБОРАТОРИЯ БИЗНЕС ЛИДЕРСТВА «ФЕНИКС»", Inches(0), H - Inches(0.65), W, Inches(0.35),
                size=10, color="5A6B73", align=PP_ALIGN.CENTER, char_spacing=1)


def build_slide_growth_rules(pres, data):
    client = data["client"]
    slide = add_slide(pres, OFFWHITE)
    title(slide, "6 базовых Правил роста")

    rules = data["growthRules"]
    col_w, col_gap, row_h, row_gap = Inches(5.55), Inches(0.3), Inches(1.1), Inches(0.25)
    start_x, start_y = Inches(0.7), Inches(1.9)
    for i, rule in enumerate(rules):
        col, row = i % 2, i // 2
        x = start_x + col * (col_w + col_gap)
        y = start_y + row * (row_h + row_gap)
        ok = rule["ok"]
        add_rounded_rect(slide, x, y, col_w, row_h,
                          fill_color=RGBColor_light_orange() if not ok else RGBColor_light_teal())
        badge_d = Inches(0.5)
        add_rounded_rect(slide, x + Inches(0.2), y + row_h // 2 - badge_d // 2, badge_d, badge_d,
                          fill_color=TERRACOTTA if not ok else TEAL, radius=0.5)
        add_textbox(slide, "!" if not ok else "✓", x + Inches(0.2), y + row_h // 2 - badge_d // 2,
                    badge_d, badge_d, size=18, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, rule["name"], x + Inches(0.9), y, col_w - Inches(1.1), row_h,
                    size=16, color=DARKTEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    footer(slide, client, dark=False)


def RGBColor_light_orange():
    from pptx.dml.color import RGBColor
    return RGBColor(0xFC, 0xE7, 0xD6)


def RGBColor_light_teal():
    from pptx.dml.color import RGBColor
    return RGBColor(0xD9, 0xEC, 0xEB)


def build_slide_agenda(pres, data):
    client = data["client"]
    slide = add_slide(pres, WHITE)
    title(slide, "Как мы проведём эти 2 часа")

    items = [
        "Сверим диагностику с Вашим собственным видением ситуации",
        "Разберём, что показали результаты — Правила роста и Вызовы",
        "Определим конкретный план действий для Вашей компании",
        "Ответим на все возникающие у Вас вопросы",
    ]
    start_y, row_h = Inches(2.0), Inches(1.15)
    for i, text in enumerate(items):
        y = start_y + i * row_h
        add_rounded_rect(slide, Inches(0.9), y, Inches(0.75), Inches(0.75), fill_color=NAVY, radius=0.5)
        add_textbox(slide, str(i + 1), Inches(0.9), y, Inches(0.75), Inches(0.75),
                    size=24, color=GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, text, Inches(2.0), y, W - Inches(3.0), Inches(0.75),
                    size=18, color=DARKTEXT, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, client, dark=False)


def build_slide_your_words(pres, data):
    client = data["client"]
    s9 = data["section9"]
    slide = add_slide(pres, OFFWHITE)
    title(slide, "Три вызова, которые Вы обозначили")
    add_textbox(slide, "Ваши собственные слова из диагностики", Inches(0.7), Inches(1.35),
                W - Inches(1.4), Inches(0.4), size=14, italic=True, color="8A7A6C")

    quotes = [s9["problem1"], s9["problem2"], s9["problem3"]]
    card_w = (W - Inches(0.7) * 2 - Inches(0.5) * 2) // 3
    card_y, card_h = Inches(2.15), Inches(4.2)
    for i, q in enumerate(quotes):
        x = Inches(0.7) + i * (card_w + Inches(0.5))
        add_rounded_rect(slide, x, card_y, card_w, card_h, fill_color=RGBColor(0xD9, 0xD9, 0xD9))
        badge_d = Inches(0.7)
        add_rounded_rect(slide, x + card_w // 2 - badge_d // 2, card_y + Inches(0.35), badge_d, badge_d,
                          fill_color=NAVY, radius=0.5)
        add_textbox(slide, str(i + 1), x + card_w // 2 - badge_d // 2, card_y + Inches(0.35), badge_d, badge_d,
                    size=22, color=GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, f"«{q}»", x + Inches(0.35), card_y + Inches(1.3), card_w - Inches(0.7), card_h - Inches(1.6),
                    size=15, italic=True, color=DARKTEXT, align=PP_ALIGN.CENTER)
    footer(slide, client, dark=False)


def build_slide_reasons(pres, data):
    client = data["client"]
    s9 = data["section9"]
    slide = add_slide(pres, OFFWHITE)
    title(slide, "Ваш взгляд на причины")

    blocks = [
        ("ВАША ВЕРСИЯ ПРИЧИНЫ", s9["whyTheseChallenges"], Inches(1.9)),
        ("ПОЧЕМУ НЕ ПОЛУЧАЕТСЯ СПРАВИТЬСЯ", s9["whyCantSolve"], Inches(4.3)),
    ]
    for label, text, y in blocks:
        add_rounded_rect(slide, Inches(0.7), y, W - Inches(1.4), Inches(2.1), fill_color=RGBColor(0xD9, 0xD9, 0xD9))
        stripe = slide.shapes.add_shape(1, Inches(0.7), y, Inches(0.12), Inches(2.1))  # 1 = RECTANGLE
        stripe.fill.solid(); stripe.fill.fore_color.rgb = TERRACOTTA; stripe.line.fill.background()
        stripe.shadow.inherit = False
        add_textbox(slide, label, Inches(1.15), y + Inches(0.25), W - Inches(2.2), Inches(0.4),
                    size=12, color="9A9088", bold=True, char_spacing=1)
        add_textbox(slide, f"«{text}»", Inches(1.15), y + Inches(0.7), W - Inches(2.2), Inches(1.25),
                    size=17, italic=True, color=DARKTEXT)
    footer(slide, client, dark=False)


def build_slide_maturity_top5(pres, data):
    from pptx.dml.color import RGBColor
    client = data["client"]
    m = data["maturityTop5"]
    slide = add_slide(pres, WHITE)
    title(slide, "Ваш индивидуальный Топ-5 вызовов")

    table_x, table_y, table_w = Inches(0.7), Inches(1.7), Inches(11.93)
    col1_w, col2_w = Inches(8.53), Inches(3.4)
    n_rows = len(m["challenges"]) + 1
    row_h = Inches(0.55)
    table_h = row_h * n_rows

    gframe = slide.shapes.add_table(n_rows, 2, table_x, table_y, table_w, table_h)
    tbl = gframe.table
    tbl.columns[0].width = col1_w
    tbl.columns[1].width = col2_w

    def _cell(cell, text, size, color, bold, bg, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(8)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT

    _cell(tbl.cell(0, 0), "Топ-5 вызовов", 13, WHITE, True, NAVY)
    _cell(tbl.cell(0, 1), "Типичная Стадия роста", 13, WHITE, True, NAVY, PP_ALIGN.CENTER)

    for i, c in enumerate(m["challenges"]):
        row = i + 1
        bg = WHITE if i % 2 == 0 else OFFWHITE
        _cell(tbl.cell(row, 0), f'{i + 1}.  «{c["text"]}»', 12.5, DARKTEXT, False, bg)
        stage_bg = TERRACOTTA if c["flagged"] else bg
        stage_color = WHITE if c["flagged"] else DARKTEXT
        _cell(tbl.cell(row, 1), c["typicalStage"], 12.5, stage_color, c["flagged"], stage_bg, PP_ALIGN.CENTER)

    plaque_y = table_y + table_h + Inches(0.35)
    plaque_h = H - plaque_y - Inches(0.6)
    add_rounded_rect(slide, table_x, plaque_y, table_w, plaque_h, fill_color=OFFWHITE, line_color=GOLD, line_width_pt=1.5)
    add_textbox(slide, "ВЫВОД", table_x + Inches(0.3), plaque_y + Inches(0.15), table_w - Inches(0.6), Inches(0.35),
                size=12, color=GOLD, bold=True)
    add_textbox(slide, m["conclusion"], table_x + Inches(0.3), plaque_y + Inches(0.55),
                table_w - Inches(0.6), plaque_h - Inches(0.7), size=14, color=DARKTEXT,
                line_spacing=1.25)

    footer(slide, client, dark=False)


def build_slide_kse_concept(pres, data):
    client = data["client"]
    slide = add_slide(pres, WHITE)
    title(slide, "11 Ключевых системных элементов")
    add_textbox(
        slide,
        "Экосистема любого устойчивого и растущего бизнеса состоит из 11 структурных "
        "элементов. Каждый устраняет корневые причины проблем — а не только их внешние "
        "симптомы — и создаёт волновой эффект позитивных изменений.",
        Inches(0.7), Inches(2.4), Inches(5.0), Inches(2.8), size=16, color=DARKTEXT,
    )
    from pptx_presentation_builder import BASE
    img_h = Inches(4.6)
    img_w = int(img_h * (1092 / 1120))
    wheel_path = BASE / "assets" / "pptx_logos" / "Колесо_элементов_с_категориями.png"
    if wheel_path.exists():
        slide.shapes.add_picture(str(wheel_path), Inches(6.3), Inches(1.5), width=img_w, height=img_h)
    footer(slide, client, dark=False)


def build_slide_symptoms(pres, data):
    from pptx.enum.shapes import MSO_SHAPE
    client = data["client"]
    slide = add_slide(pres, WHITE)
    title(slide, "Симптомы и корневые причины")

    # Раньше нижняя голубая «вода» доходила почти до самого низа слайда и
    # наезжала на логотип в футере — весь блок (кроме футера) поднят ближе
    # к заголовку, а вода останавливается с явным отступом от футера.
    tip_y = Inches(1.3)
    mid_y = Inches(3.5)
    water_bottom = H - Inches(0.75)
    ice_x = W // 2

    water = slide.shapes.add_shape(1, 0, mid_y, W, water_bottom - mid_y)
    water.fill.solid(); water.fill.fore_color.rgb = RGBColor(0xDC, 0xEA, 0xF0); water.line.fill.background()
    water.shadow.inherit = False

    tip = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, ice_x - Inches(1.6), tip_y, Inches(3.2), mid_y - tip_y)
    tip.fill.solid(); tip.fill.fore_color.rgb = TERRACOTTA; tip.line.fill.background(); tip.shadow.inherit = False
    tip.rotation = 180
    add_textbox(slide, "ВИДНО СРАЗУ", Inches(2.77), tip_y + Inches(0.1), Inches(2.0), Inches(0.35),
                size=11, color=TERRACOTTA, bold=True, align=PP_ALIGN.RIGHT)
    add_textbox(slide, ["Правила роста, которые не соблюдаются", "", "Вызовы, с которыми сталкивается бизнес"],
                Inches(1.77), tip_y + Inches(0.55), Inches(3.0), Inches(1.5), size=12.5, color=DARKTEXT, align=PP_ALIGN.RIGHT)

    base_tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, ice_x - Inches(2.6), mid_y, Inches(5.2), Inches(2.2))
    base_tri.fill.solid(); base_tri.fill.fore_color.rgb = NAVY; base_tri.line.fill.background(); base_tri.shadow.inherit = False
    add_textbox(slide, "ГЛУБЖЕ — КОРЕНЬ", Inches(9.57), mid_y + Inches(0.3), Inches(2.5), Inches(0.35),
                size=11, color=GOLD, bold=True)
    add_textbox(slide, "Отсутствующие Ключевые системные элементы", Inches(9.57), mid_y + Inches(0.75),
                Inches(3.0), Inches(1.0), size=12.5, color=DARKTEXT)

    add_textbox(
        slide,
        "Правила роста и Вызовы — это внешние симптомы более глубоких проблем. Именно "
        "корневые причины устраняются путём внедрения недостающих Ключевых системных элементов.",
        Inches(0.9), H - Inches(1.35), W - Inches(1.8), Inches(0.55), size=13, italic=True,
        color="5A6B73", align=PP_ALIGN.CENTER,
    )
    footer(slide, client, dark=False)


def build_slide_priority_chain(pres, data):
    """Слайд №9 — полностью динамический: и левая колонка (Правила роста /
    Вызовы / Области — до 6 строк каждая), и правая колонка (КСЭ, к которым
    привела диагностика — от 1 до 10+ плашек) подстраивают высоту и размер
    шрифта под реальное количество пунктов конкретного клиента."""
    client = data["client"]
    ec = data["evidenceChain"]
    slide = add_slide(pres, WHITE)
    title(slide, "Как мы вышли на эти Элементы")

    left_x, left_w = Inches(0.6), Inches(5.3)
    right_x, right_w = Inches(8.9), Inches(3.8)
    arrow_x1, arrow_x2 = left_x + left_w + Inches(0.15), right_x - Inches(0.15)

    top_y = Inches(1.5)
    bottom_y = H - Inches(0.75)
    avail_h = bottom_y - top_y

    # --- левая колонка: 3 плашки-источника, высота каждой зависит от числа
    # пунктов внутри (до 6 строк на плашку) ---
    ev_groups = [
        ("ПРАВИЛА РОСТА (не выполняются)", ec.get("growthRuleMismatches") or []),
        ("ВЫЗОВЫ (Вы отметили как острые: 6+ баллов)", ec.get("challenges") or []),
        ("ОБЛАСТИ БИЗНЕСА (просели)", ec.get("areas") or []),
    ]
    max_items = max((len(items) for _, items in ev_groups), default=1) or 1
    header_h = Inches(0.38)
    left_gap_base = Inches(0.16)

    def _compute_heights(bullet_sz):
        lh = int(Inches(bullet_sz / 72 * 1.55))
        hs = [header_h + Inches(0.12) + lh * max(len(items), 1) + Inches(0.12) for _, items in ev_groups]
        return hs, lh

    bullet_size = 11.5 if max_items <= 4 else max(8.0, 11.5 - (max_items - 4) * 0.55)
    heights, line_h = _compute_heights(bullet_size)
    left_gap = left_gap_base
    total_h = sum(heights) + left_gap * (len(ev_groups) - 1)
    # Сначала жмём шрифт (до 7pt) — тогда текст и плашка остаются в
    # согласии друг с другом. Плашки целиком сжимаем только в
    # экстремальном случае, когда даже 7pt не помогает.
    while total_h > avail_h and bullet_size > 7.0:
        bullet_size -= 0.5
        heights, line_h = _compute_heights(bullet_size)
        total_h = sum(heights) + left_gap * (len(ev_groups) - 1)
    if total_h > avail_h:
        scale = avail_h / total_h
        heights = [int(h * scale) for h in heights]
        left_gap = int(left_gap * scale)

    y = top_y
    box_positions = []
    for (label, items), h in zip(ev_groups, heights):
        add_rounded_rect(slide, left_x, y, left_w, h, fill_color=OFFWHITE, line_color=TERRACOTTA, line_width_pt=1.5, radius=0.06)
        add_textbox(slide, label, left_x + Inches(0.25), y + Inches(0.1), left_w - Inches(0.5), header_h,
                    size=10.5, color=TERRACOTTA, bold=True)
        # 22.08.2026: блок "ВЫЗОВЫ" — нумерованный список (не буллеты), а
        # два других блока ("Правила роста", "Области") — как были, с
        # буллетами. Номера нужны, чтобы под каждым КСЭ в правой колонке
        # можно было компактно сослаться "Вызов 1, 7" вместо полного текста
        # (см. kseCaptions в pptx_data_export.py) — нумерация здесь и там
        # обязана совпадать: и тут, и в _build_evidence_chain порядок
        # пунктов ОДИНАКОВЫЙ (одна и та же сортировка по убыванию балла),
        # так что просто enumerate() по items даёт тот же номер, что уже
        # использован при построении kseCaptions.
        if label.startswith("ВЫЗОВЫ") and items:
            numbered_items = [f"{i + 1}. {item}" for i, item in enumerate(items)]
            add_textbox(slide, numbered_items, left_x + Inches(0.25), y + header_h, left_w - Inches(0.5), h - header_h - Inches(0.08),
                        size=bullet_size, color=DARKTEXT, line_spacing=1.1)
        else:
            add_bullets(slide, items or ["—"], left_x + Inches(0.25), y + header_h, left_w - Inches(0.5), h - header_h - Inches(0.08),
                        size=bullet_size, color=DARKTEXT)
        box_positions.append((y, h))
        y += h + left_gap

    # --- правая колонка: реальные КСЭ, к которым привела эта цепочка —
    # от 1 до 10+ плашек, высота и шрифт подстраиваются под количество ---
    # 22.08.2026: показываем ВСЕ рекомендованные КСЭ (не только первые 3) —
    # Игорь подтвердил, что урезка top-3 не нужна.
    related_kse = ec.get("relatedKse") or data["allKseOrdered"]
    n_kse = max(len(related_kse), 1)
    kse_gap = min(Inches(0.2), int(avail_h / n_kse * 0.18))
    kse_h = min(Inches(1.15), (avail_h - kse_gap * (n_kse - 1)) // n_kse)
    kse_font = 13 if n_kse <= 3 else max(9.0, 13 - (n_kse - 3) * 0.7)
    # Подпись-обоснование под названием — компактная ссылка на конкретные
    # пункты левой колонки ("Вызов 1, 7" и т.п., см. kseCaptions в
    # pptx_data_export.py). Заметно мельче названия, чтобы не спорить с ним
    # за внимание, но всегда читаема (не меньше 7pt).
    caption_font = max(7.0, kse_font - 4)
    kse_captions = ec.get("kseCaptions") or {}

    ry = top_y
    right_bottom = ry
    for name in related_kse:
        add_rounded_rect(slide, right_x, ry, right_w, kse_h, fill_color=NAVY)
        caption = kse_captions.get(name, "")
        if caption:
            # Название — верхняя часть плашки (anchor BOTTOM, чтобы текст
            # "прижался" к линии раздела с подписью, а не висел в воздухе
            # посередине пустой верхней половины).
            add_textbox(slide, name, right_x + Inches(0.2), ry + Inches(0.05), right_w - Inches(0.4), kse_h * 0.6,
                        size=kse_font, color=GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
            add_textbox(slide, caption, right_x + Inches(0.2), ry + kse_h * 0.6, right_w - Inches(0.4), kse_h * 0.4 - Inches(0.05),
                        size=caption_font, color=RGBColor(0xB9, 0xAF, 0xA5), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        else:
            add_textbox(slide, name, right_x + Inches(0.2), ry, right_w - Inches(0.4), kse_h,
                        size=kse_font, color=GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        ry += kse_h + kse_gap
        right_bottom = ry - kse_gap

    right_center_y = top_y + (right_bottom - top_y) // 2
    for by, bh in box_positions:
        add_line_arrow(slide, arrow_x1, by + bh // 2, arrow_x2, right_center_y, RGBColor(0xB9, 0xAF, 0xA5), 1.75)

    footer(slide, client, dark=False)


def build_slide_cost_only(pres, data):
    client = data["client"]
    s9 = data["section9"]
    slide = add_slide(pres, NAVY)
    add_textbox(slide, "ЧЕРЕЗ 12 МЕСЯЦЕВ, ЕСЛИ НИЧЕГО НЕ ИЗМЕНИТСЯ", Inches(1.2), Inches(1.2), W - Inches(2.4), Inches(0.8),
                size=24, color=TERRACOTTA, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, f'«{s9["costOfInaction"]}»', Inches(1.6), Inches(2.6), W - Inches(3.2), Inches(1.6),
                size=20, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, f'«{s9["priceOfInaction"]}»', Inches(1.6), Inches(4.4), W - Inches(3.2), Inches(2.0),
                size=17, italic=True, color="C9D1D5", align=PP_ALIGN.CENTER)
    footer(slide, client, dark=True)


def build_slide_full_contrast(pres, data):
    client = data["client"]
    s9 = data["section9"]
    slide = add_slide(pres, WHITE)
    half_w = W // 2

    left = slide.shapes.add_shape(1, 0, 0, half_w, H)
    left.fill.solid(); left.fill.fore_color.rgb = NAVY; left.line.fill.background(); left.shadow.inherit = False
    add_textbox(slide, "ЧЕРЕЗ 12 МЕСЯЦЕВ, ЕСЛИ НИЧЕГО НЕ ИЗМЕНИТСЯ", Inches(0.6), Inches(0.6), half_w - Inches(1.2), Inches(0.8),
                size=17, color=TERRACOTTA, bold=True)
    add_textbox(slide, f'«{s9["costOfInaction"]}»', Inches(0.6), Inches(1.7), half_w - Inches(1.2), Inches(1.6),
                size=15, italic=True, color=WHITE)
    add_textbox(slide, f'«{s9["priceOfInaction"]}»', Inches(0.6), Inches(3.6), half_w - Inches(1.2), Inches(2.2),
                size=13.5, italic=True, color="C9D1D5")

    right = slide.shapes.add_shape(1, half_w, 0, half_w, H)
    right.fill.solid(); right.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xEC); right.line.fill.background(); right.shadow.inherit = False
    add_textbox(slide, "ЕСЛИ БЫ СБЫЛАСЬ ВАША МЕЧТА", half_w + Inches(0.6), Inches(0.6), half_w - Inches(1.2), Inches(0.8),
                size=17, color=NAVY, bold=True)
    add_textbox(slide, f'«{s9["dreamOutcome"]}»', half_w + Inches(0.6), Inches(1.7), half_w - Inches(1.2), Inches(3.5),
                size=16, italic=True, color=DARKTEXT)

    add_rounded_rect(slide, half_w - Inches(0.45), H // 2 - Inches(0.45), Inches(0.9), Inches(0.9),
                      fill_color=GOLD, line_color=WHITE, line_width_pt=3, radius=0.5)
    add_textbox(slide, "VS", half_w - Inches(0.45), H // 2 - Inches(0.45), Inches(0.9), Inches(0.9),
                size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, client, dark=False)


def build_slide_how_to_reach(pres, data):
    client = data["client"]
    slide = add_slide(pres, NAVY)
    add_textbox(slide, "Как выйти на этот результат?", Inches(1.2), Inches(2.6), W - Inches(2.4), Inches(0.9),
                size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Что нужно сделать системно, шаг за шагом, чтобы прийти к тому будущему, которое Вы описали?",
                Inches(2.0), Inches(3.7), W - Inches(4.0), Inches(1.2), size=18, color=GOLD, align=PP_ALIGN.CENTER)
    footer(slide, client, dark=True)


def build_slide_kse_list_repeat(pres, data):
    """Слайд №13 — ВЕСЬ список рекомендованных клиенту КСЭ (не top-N обрезка):
    это важно, т.к. количество Программ на слайдах №20-23 («Дом» / карточки
    по ярусам) должно 1-в-1 совпадать с этим списком — оба берутся из
    allKseOrdered. На Стадиях 6-7 в Фундаменте пунктов может быть до 10."""
    client = data["client"]
    priority = data.get("allKseOrdered") or []
    slide = add_slide(pres, RGBColor(0xF7, 0xF5, 0xF2))
    title(slide, "Результаты диагностики")
    add_textbox(slide, "Необходимо внедрение следующих Ключевых системных элементов:",
                Inches(0.7), Inches(1.4), W - Inches(1.4), Inches(0.5), size=15, color=DARKTEXT)

    start_y, end_y = Inches(2.2), H - Inches(0.9)
    avail_h = end_y - start_y
    n = max(len(priority), 1)
    gap = min(Inches(0.3), int(avail_h / n * 0.22))
    card_h = (avail_h - gap * (n - 1)) // n
    # Шрифт названия и кружка-номера плавно уменьшаются с ростом количества
    # плашек, чтобы 10 пунктов (Стадии 6-7, Фундамент) не схлопывались.
    name_size = max(11.0, min(18.0, 18.0 - (n - 5) * 0.9)) if n > 5 else 18.0
    for i, name in enumerate(priority):
        y = start_y + i * (card_h + gap)
        add_rounded_rect(slide, Inches(1.5), y, W - Inches(3.0), card_h, fill_color=NAVY)
        circle_d = min(Inches(0.64), int(card_h * 0.7))
        add_rounded_rect(slide, Inches(1.85), y + card_h // 2 - circle_d // 2, circle_d, circle_d,
                          fill_color=TERRACOTTA, radius=0.5)
        circle_d_in = circle_d / 914400
        add_textbox(slide, str(i + 1), Inches(1.85), y + card_h // 2 - circle_d // 2, circle_d, circle_d,
                    size=min(22, circle_d_in * 32), color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, name, Inches(2.8), y, W - Inches(5.0), card_h,
                    size=name_size, color=GOLD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, client, dark=False)


def build_slide_what_order(pres, data):
    client = data["client"]
    slide = add_slide(pres, NAVY)
    add_textbox(slide, "В каком порядке внедрять Ключевые системные элементы?",
                Inches(1.2), Inches(3.0), W - Inches(2.4), Inches(1.4),
                size=28, color=GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, client, dark=True)


_HOUSE_TIERS = [
    ("НАДСТРОЙКА", "Надстройка", "3", RGBColor(0xCB, 0xB7, 0x9A), NAVY, Inches(1.85), Inches(1.4)),
    ("ЯДРО", "Ядро", "2", TEAL, WHITE, Inches(3.45), Inches(1.4)),
    ("ФУНДАМЕНТ", "Фундамент", "1", NAVY, WHITE, Inches(5.05), Inches(1.4)),
]


def _draw_house(slide, note_content_fn):
    """note_content_fn(tier_name) -> str — текст правой заметки для яруса."""
    house_x, house_w = Inches(0.9), Inches(5.6)
    for key, tier_name, num, color, text_color, y, h in _HOUSE_TIERS:
        block = slide.shapes.add_shape(1, house_x, y, house_w, h)
        block.fill.solid(); block.fill.fore_color.rgb = color
        block.line.color.rgb = WHITE; block.line.width = Pt(2); block.shadow.inherit = False
        badge_d = Inches(0.8)
        add_rounded_rect(slide, house_x + Inches(0.3), y + h // 2 - badge_d // 2, badge_d, badge_d,
                          fill_color=TERRACOTTA, line_color=WHITE, line_width_pt=2, radius=0.5)
        add_textbox(slide, num, house_x + Inches(0.3), y + h // 2 - badge_d // 2, badge_d, badge_d,
                    size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, key, house_x + Inches(1.3), y + Inches(0.2), house_w - Inches(1.5), Inches(0.4),
                    size=16, color=text_color, bold=True)

    note_x = house_x + house_w + Inches(0.6)
    note_w = W - note_x - Inches(0.6)
    for key, tier_name, num, color, text_color, y, h in _HOUSE_TIERS:
        add_rounded_rect(slide, note_x, y, note_w, h, fill_color=OFFWHITE, radius=0.06)
        content, bold = note_content_fn(tier_name)
        add_textbox(slide, content, note_x + Inches(0.25), y, note_w - Inches(0.5), h,
                    size=12.5 if bold else 12, color=DARKTEXT, bold=bold, anchor=MSO_ANCHOR.MIDDLE)

    add_textbox(
        slide,
        "Попытка внедрить более поздний уровень раньше предыдущего не ускоряет результат, "
        "а создаёт видимость порядка без реальной устойчивости.",
        Inches(0.9), Inches(6.6), W - Inches(1.8), Inches(0.55), size=11.5, italic=True,
        color="6B6259", align=PP_ALIGN.CENTER,
    )


_TIER_DESCRIPTIONS = {
    "Надстройка": "Достраивается, когда основа уже устойчива",
    "Ядро": "Несущие стены — начинаются, когда Фундамент выдержит нагрузку",
    "Фундамент": "То, без чего невозможно опираться на остальное — закладывается первым",
}


def build_slide_house(pres, data):
    client = data["client"]
    slide = add_slide(pres, WHITE)
    title(slide, "Порядок внедрения")
    _draw_house(slide, lambda tier_name: (_TIER_DESCRIPTIONS[tier_name], False))
    footer(slide, client, dark=False)


def build_slide_house_continuation(pres, data):
    """Слайд №20 — в отличие от build_slide_house (общая, концептуальная
    версия «Дома» с короткими фиксированными описаниями ярусов), здесь
    названия Программ клиент-специфичны и их может быть до 10 на ярус
    (Фундамент на Стадиях 6-7). Поэтому раскладка отличается: плашки ярусов
    на всю ширину слайда, названия Программ — в 2 колонки, высота плашки
    динамически подстраивается под число Программ в ней."""
    client = data["client"]
    tier_programs = data["tierPrograms"]
    slide = add_slide(pres, WHITE)
    title(slide, "Порядок внедрения Программ")

    tiers = [
        ("НАДСТРОЙКА", "3", RGBColor(0xCB, 0xB7, 0x9A), NAVY, tier_programs.get("Надстройка", [])),
        ("ЯДРО", "2", TEAL, WHITE, tier_programs.get("Ядро", [])),
        ("ФУНДАМЕНТ", "1", NAVY, WHITE, tier_programs.get("Фундамент", [])),
    ]

    house_x, house_w = Inches(0.7), W - Inches(1.4)
    top_y = Inches(1.5)
    bottom_y = H - Inches(0.85)
    avail_h = bottom_y - top_y
    gap = Inches(0.16)

    header_h = Inches(0.42)
    row_h_base = Inches(0.27)
    min_tier_h = Inches(0.75)

    heights = []
    for _, _, _, _, programs in tiers:
        n = len(programs)
        rows = max(1, -(-n // 2)) if n else 1  # ceil(n / 2 колонки)
        h = max(min_tier_h, header_h + rows * row_h_base + Inches(0.22))
        heights.append(h)

    total_h = sum(heights) + gap * (len(tiers) - 1)
    if total_h > avail_h:
        scale = avail_h / total_h
        heights = [int(h * scale) for h in heights]
        gap = int(gap * scale)
        row_h_base = int(row_h_base * scale)

    y = top_y
    for (key, num, color, text_color, programs), h in zip(tiers, heights):
        block = slide.shapes.add_shape(1, house_x, y, house_w, h)
        block.fill.solid(); block.fill.fore_color.rgb = color
        block.line.color.rgb = WHITE; block.line.width = Pt(2); block.shadow.inherit = False

        badge_d = min(Inches(0.7), h - Inches(0.15))
        add_rounded_rect(slide, house_x + Inches(0.25), y + Inches(0.12), badge_d, badge_d,
                          fill_color=TERRACOTTA, line_color=WHITE, line_width_pt=2, radius=0.5)
        add_textbox(slide, num, house_x + Inches(0.25), y + Inches(0.12), badge_d, badge_d,
                    size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, key, house_x + Inches(1.15), y + Inches(0.1), Inches(3.0), Inches(0.32),
                    size=15, color=text_color, bold=True)

        content_y = y + header_h
        content_h = h - header_h - Inches(0.1)
        content_x = house_x + Inches(1.15)
        content_w = house_w - Inches(1.45)
        if not programs:
            add_textbox(slide, "— не требуется для Вашего плана", content_x, content_y,
                        content_w, content_h, size=12, color=text_color)
        else:
            n = len(programs)
            names = [p["name"] for p in programs]
            col_n = max(1, -(-n // 2))
            col1, col2 = names[:col_n], names[col_n:]
            font_size = 12.0 if n <= 6 else max(9.0, 12.0 - (n - 6) * 0.5)
            col_w = (content_w - Inches(0.3)) // 2
            add_bullets(slide, col1, content_x, content_y, col_w, content_h, size=font_size, color=text_color)
            if col2:
                add_bullets(slide, col2, content_x + col_w + Inches(0.3), content_y, col_w, content_h,
                            size=font_size, color=text_color)
        y += h + gap

    footer(slide, client, dark=False)


def build_slide_how_to_implement(pres, data):
    client = data["client"]
    slide = add_slide(pres, WHITE)
    title(slide, "Как внедрять Ключевые системные элементы?")

    box_y, box_h, box_w = Inches(3.1), Inches(1.5), Inches(3.6)
    x1, x2 = W // 2 - box_w - Inches(1.1), W // 2 + Inches(1.1)

    add_rounded_rect(slide, x1, box_y, box_w, box_h, fill_color=NAVY)
    add_textbox(slide, ["1 Ключевой", "системный элемент"], x1, box_y, box_w, box_h,
                size=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rounded_rect(slide, x2, box_y, box_w, box_h, fill_color=TEAL)
    add_textbox(slide, ["1 Консалтинговая", "программа"], x2, box_y, box_w, box_h,
                size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_line_arrow(slide, x1 + box_w + Inches(0.15), box_y + box_h // 2,
                    x2 - Inches(0.15), box_y + box_h // 2, TERRACOTTA, 3)

    add_textbox(
        slide, "Внедрение Ключевого системного элемента происходит в виде одноимённой Консалтинговой программы.",
        Inches(1.5), Inches(5.2), W - Inches(3.0), Inches(0.8), size=15, italic=True,
        color=DARKTEXT, align=PP_ALIGN.CENTER,
    )
    footer(slide, client, dark=False)


def build_slide_what_is_program(pres, data):
    client = data["client"]
    slide = add_slide(pres, WHITE)
    title(slide, "Что такое Консалтинговая программа?")
    add_textbox(
        slide,
        "Это структурированный процесс внедрения одного Ключевого системного "
        "элемента — с чёткими этапами, инструментами и результатами на каждом шаге.",
        Inches(0.7), Inches(2.4), Inches(4.8), Inches(1.9), size=18, color=NAVY, bold=True,
    )
    ln = slide.shapes.add_connector(1, Inches(0.7), Inches(4.5), Inches(5.2), Inches(4.5))
    ln.line.color.rgb = TEAL; ln.line.width = Pt(1.5)
    add_textbox(
        slide,
        "Не разовая консультация, а последовательная работа на протяжении недель "
        "или месяцев, пока Элемент не станет частью того, как устроен Ваш бизнес.",
        Inches(0.7), Inches(4.7), Inches(4.8), Inches(1.6), size=14.5, color=DARKTEXT,
    )
    img_path = BASE / "assets" / "pptx_logos" / "Рисунок3.png"
    if img_path.exists():
        img_h = Inches(4.3)
        img_w = int(img_h * (2207 / 1677))
        slide.shapes.add_picture(str(img_path), Inches(6.1), Inches(2.2), width=img_w, height=img_h)
    footer(slide, client, dark=False)


def build_slide_program_model(pres, data):
    client = data["client"]
    slide = add_slide(pres, WHITE)
    title(slide, "Модель передачи знаний и отработки навыков")
    add_textbox(
        slide,
        "Авторская «Модель передачи знаний и отработки навыков их применения» — "
        "основана на теории обучения взрослых, использует мультимодальные методы. "
        "Это способствует быстрому усвоению и применению на практике.",
        Inches(0.7), Inches(1.5), W - Inches(1.4), Inches(1.1), size=15, color=DARKTEXT,
    )
    elements = ["Обучающие видео и текстовые материалы", "Рабочие сессии", "Практическое применение"]
    colors = [NAVY, TEAL, TERRACOTTA]
    card_w = (W - Inches(0.7) * 2 - Inches(0.4) * 2) // 3
    card_y, card_h = Inches(3.0), Inches(3.2)
    for i, label in enumerate(elements):
        x = Inches(0.7) + i * (card_w + Inches(0.4))
        add_rounded_rect(slide, x, card_y, card_w, card_h, fill_color=RGBColor(0xD9, 0xD9, 0xD9))
        badge_d = Inches(1.0)
        add_rounded_rect(slide, x + card_w // 2 - badge_d // 2, card_y + Inches(0.4), badge_d, badge_d,
                          fill_color=colors[i], radius=0.5)
        add_textbox(slide, str(i + 1), x + card_w // 2 - badge_d // 2, card_y + Inches(0.4), badge_d, badge_d,
                    size=30, color=GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, label, x + Inches(0.3), card_y + Inches(1.65), card_w - Inches(0.6), Inches(1.4),
                    size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, client, dark=False)


def build_slide_program_model_practice(pres, data):
    client = data["client"]
    slide = add_slide(pres, WHITE)
    title(slide, "Модель на практике")

    steps = [
        ("Изучение теоретических принципов", "В LMS-системе",
         ["Обучающие видео", "Схемы", "Таблицы", "Гиды"], NAVY),
        ("Применение теории для разработки решений", "В Zoom",
         ["Индивидуальная и командная работа на Рабочих сессиях", "Разработка готовых решений"], TEAL),
        ("Применение разработанных решений", "В компании",
         ["Апробация решений", "Применение в рабочем процессе для решения конкретных проблем бизнеса"], TERRACOTTA),
    ]
    card_w, gap, card_y, card_h = Inches(3.5), Inches(0.55), Inches(2.0), Inches(4.1)
    total_w = card_w * 3 + gap * 2
    start_x = (W - total_w) // 2

    for i, (step_title, where, items, color) in enumerate(steps):
        x = start_x + i * (card_w + gap)
        add_rounded_rect(slide, x, card_y, card_w, card_h, fill_color=RGBColor(0xD9, 0xD9, 0xD9))
        badge_d = Inches(0.7)
        add_rounded_rect(slide, x + card_w // 2 - badge_d // 2, card_y + Inches(0.3), badge_d, badge_d,
                          fill_color=color, radius=0.5)
        add_textbox(slide, str(i + 1), x + card_w // 2 - badge_d // 2, card_y + Inches(0.3), badge_d, badge_d,
                    size=22, color=GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, step_title, x + Inches(0.25), card_y + Inches(1.15), card_w - Inches(0.5), Inches(0.9),
                    size=13.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, where.upper(), x + Inches(0.25), card_y + Inches(2.05), card_w - Inches(0.5), Inches(0.3),
                    size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_bullets(slide, items, x + Inches(0.3), card_y + Inches(2.45), card_w - Inches(0.6), card_h - Inches(2.6),
                    size=11, color=DARKTEXT)
        if i < 2:
            add_line_arrow(slide, x + card_w + Inches(0.05), card_y + card_h // 2,
                            x + card_w + gap - Inches(0.05), card_y + card_h // 2, TERRACOTTA, 6)

    loop_y = card_y + card_h + Inches(0.3)
    mid_h = card_y + card_h // 2
    right_edge = start_x + 2 * (card_w + gap) + card_w
    left_edge = start_x
    stub = Inches(0.35)
    line_color = TERRACOTTA

    def straight(x1, y1, x2, y2, arrow=False):
        c = slide.shapes.add_connector(1, x1, y1, x2, y2)
        c.line.color.rgb = line_color
        c.line.width = Pt(6)
        if arrow:
            from pptx.oxml.ns import qn
            ln = c.line._get_or_add_ln()
            tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle"})
            ln.append(tail)

    straight(right_edge, mid_h, right_edge + stub, mid_h)
    straight(right_edge + stub, mid_h, right_edge + stub, loop_y)
    straight(left_edge - stub, loop_y, right_edge + stub, loop_y)
    straight(left_edge - stub, mid_h, left_edge - stub, loop_y)
    straight(left_edge - stub, mid_h, left_edge, mid_h, arrow=True)

    add_textbox(slide, "ПОВТОРЯЮЩИЙСЯ ПРОЦЕСС", start_x, loop_y + Inches(0.1), total_w, Inches(0.25),
                size=10, color=TERRACOTTA, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, client, dark=False)


# Справочник ссылок на Калькуляторы окупаемости — НЕ клиент-специфичные
# данные, общий для всех клиентов (как и в base.js) — поэтому не приходит
# из pptx_data.json. Источник ссылок: файл «Ссылки на калькуляторы.xlsx».
ROI_CALCULATORS = {
    "Сильная Управленческая команда": "https://fenix-lab.ru/smt-prices#ROIcalculator",
    "Ключевые показатели эффективности": "https://fenix-lab.ru/kpi-program#ROIcalculator",
    "Комплексное планирование": "https://fenix-lab.ru/kp#ROIcalculator",
    "Ценности бренда и Базовые ценности": "https://fenix-lab.ru/values-program#ROIcalculator",
    "Критерии роста бизнеса": "https://fenix-lab.ru/krb-program#ROIcalculator",
    "Базовые бизнес-процессы": "https://fenix-lab.ru/bbp#ROIcalculator",
    "Структура рабочих совещаний": "https://fenix-lab.ru/meetings-program#ROIcalculator",
    "Структура Развития бизнеса": "https://fenix-lab.ru/srb-program#ROIcalculator",
    "Коучинговое управление персоналом": "https://fenix-lab.ru/kup#ROIcalculator",
    "Организационная структура": "https://fenix-lab.ru/org-structura#ROIcalculator",
    "Бизнес-модель": "https://fenix-lab.ru/bm#ROIcalculator",
    "Возрождение малого бизнеса": "https://fenix-lab.ru/vmb#ROIcalculator",
}


MAX_CARDS_PER_SLIDE = 3  # больше — карточки становятся нечитаемыми (текст
                          # схлопывается в один символ на строку); лучше
                          # несколько слайдов подряд, чем один нечитаемый.


def _build_program_cards_slide(pres, data, tier_name, slide_title):
    client = data["client"]
    programs = data["tierPrograms"].get(tier_name, [])
    if not programs:
        slide = add_slide(pres, WHITE)
        title(slide, slide_title)
        add_textbox(slide, "— не требуется для Вашего плана", Inches(0.7), Inches(2.5), W - Inches(1.4), Inches(0.6),
                    size=16, color=DARKTEXT)
        footer(slide, client, dark=False)
        return

    chunks = [programs[i:i + MAX_CARDS_PER_SLIDE] for i in range(0, len(programs), MAX_CARDS_PER_SLIDE)]
    for page_i, chunk in enumerate(chunks):
        slide = add_slide(pres, WHITE)
        page_title = slide_title if len(chunks) == 1 else f"{slide_title} ({page_i + 1}/{len(chunks)})"
        title(slide, page_title)

        card_w = (W - Inches(0.7) * 2 - Inches(0.4) * (len(chunk) - 1)) // len(chunk)
        card_y, card_h = Inches(1.9), Inches(4.9)
        for i, p in enumerate(chunk):
            x = Inches(0.7) + i * (card_w + Inches(0.4))

            # 22.08.2026: строка "Калькулятор окупаемости инвестиций: URL"
            # раньше стояла ВНУТРИ карточки (над разделительной линией) и на
            # части Программ визуально наезжала на 3-й буллет блока "Что Вы
            # получаете" или на саму разделительную линию (длинные названия
            # Программ/буллеты не оставляли ей запаса по высоте). Убрана из
            # карточки полностью и вынесена НАД карточкой, в свободные ~0.5"
            # между заголовком слайда и верхним краем карточки — с этим
            # решён и визуальный overflow, и сама ссылка стала кликабельной
            # (без видимого URL в тексте — только слова "Калькулятор
            # окупаемости инвестиций", per требование Игоря 22.08.2026:
            # Презентация не выдаётся клиенту файлом, ссылка открывается
            # Игорем прямо во время трансляции по просьбе клиента).
            calc_url = ROI_CALCULATORS.get(p["name"])
            if calc_url:
                calc_box = add_textbox(slide, "🔗 Калькулятор окупаемости инвестиций",
                                        x + Inches(0.05), card_y - Inches(0.34), card_w - Inches(0.1), Inches(0.3),
                                        size=11, color=TEAL, bold=False)
                calc_run = calc_box.text_frame.paragraphs[0].runs[0]
                calc_run.hyperlink.address = calc_url
                calc_run.font.underline = True

            add_rounded_rect(slide, x, card_y, card_w, card_h, fill_color=RGBColor(0xDE, 0xEB, 0xF7),
                              line_color=RGBColor(0xE4, 0xE1, 0xDC), line_width_pt=1)
            stripe = slide.shapes.add_shape(1, x, card_y, card_w, Inches(0.1))
            stripe.fill.solid(); stripe.fill.fore_color.rgb = NAVY; stripe.line.fill.background(); stripe.shadow.inherit = False

            add_textbox(slide, p["name"], x + Inches(0.3), card_y + Inches(0.3), card_w - Inches(0.6), Inches(0.9),
                        size=16, color=NAVY, bold=True)
            add_textbox(slide, "ЧТО ВЫ ПОЛУЧАЕТЕ:", x + Inches(0.3), card_y + Inches(1.3), card_w - Inches(0.6), Inches(0.3),
                        size=10, color="9A9088", bold=True)
            # Строка калькулятора больше не занимает место ВНУТРИ карточки —
            # блок буллетов теперь всегда получает полную высоту 2.2",
            # независимо от того, есть ли у Программы калькулятор.
            add_bullets(slide, p["outcomes"], x + Inches(0.3), card_y + Inches(1.65), card_w - Inches(0.6), Inches(2.2),
                        size=11.5, color=DARKTEXT)

            ln = slide.shapes.add_connector(1, x + Inches(0.3), card_y + Inches(3.95), x + card_w - Inches(0.3), card_y + Inches(3.95))
            ln.line.color.rgb = RGBColor(0xD8, 0xD3, 0xCC); ln.line.width = Pt(1)
            add_textbox(slide, f'Срок: {p["duration"]}', x + Inches(0.3), card_y + Inches(4.1), card_w - Inches(0.6), Inches(0.3),
                        size=12, color=DARKTEXT)
            add_textbox(slide, p["price"], x + Inches(0.3), card_y + Inches(4.4), card_w - Inches(0.6), Inches(0.4),
                        size=20, color=TERRACOTTA, bold=True)
        footer(slide, client, dark=False)


def build_slide_programs_foundation(pres, data):
    _build_program_cards_slide(pres, data, "Фундамент", "Фундамент: Программы")


def build_slide_programs_core(pres, data):
    _build_program_cards_slide(pres, data, "Ядро", "Ядро: Программы")


def build_slide_programs_superstructure(pres, data):
    _build_program_cards_slide(pres, data, "Надстройка", "Надстройка: Программы")


def _fmt_rub(n):
    return f'{n:,.0f}'.replace(',', ' ') + ' ₽'


def build_slide_guarantee(pres, data):
    client = data["client"]
    slide = add_slide(pres, NAVY)
    badge_size = Inches(2.2)
    badge_path = BASE / "assets" / "pptx_logos" / "Знак_Гарантия_качества_transparent.png"
    if badge_path.exists():
        slide.shapes.add_picture(str(badge_path), W // 2 - badge_size // 2, Inches(0.75), width=badge_size, height=badge_size)
    add_textbox(slide, "СТРАХОВКА ОТ СОМНЕНИЙ", Inches(1.5), Inches(3.35), W - Inches(3), Inches(0.6),
                size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(
        slide,
        "Мы полностью уверены в результате. Если после первой Рабочей сессии Программы Вы "
        "не увидите реальной пользы — Вы получаете 100% возврат средств без каких-либо "
        "дополнительных условий.",
        Inches(2.3), Inches(4.15), W - Inches(4.6), Inches(1.6), size=16, color="C9D1D5", align=PP_ALIGN.CENTER,
    )
    footer(slide, client, dark=True)


def build_slide_loyalty_program(pres, data):
    from pptx.oxml.ns import qn
    client = data["client"]
    items = data["loyaltyProgram"]["items"]
    slide = add_slide(pres, WHITE)
    title(slide, "Программа лояльности")
    add_textbox(
        slide,
        "Скидка не зависит от порядка внедрения (тот определяется только структурой «Дома») "
        "— это отдельное правило, которое действует поверх него: каждая следующая Программа "
        "в Вашей последовательности стоит дешевле предыдущей.",
        Inches(0.7), Inches(1.15), Inches(11.93), Inches(0.5), size=11, italic=True, color=DARKTEXT,
    )

    table_x, table_y, table_w = Inches(0.7), Inches(1.75), Inches(11.93)
    col_w = [Inches(0.55), Inches(4.6), Inches(2.2), Inches(1.3), Inches(3.28)]
    n_rows = len(items) + 1

    # Высота строки и таблицы подстраиваются под количество Программ (до 11
    # на Стадиях 6-7), чтобы плашка «Ваша экономия…» не наезжала на футер.
    footer_top = H - Inches(0.75)
    plaque_h = Inches(0.55)
    gap_before_plaque = Inches(0.2)
    max_table_h = footer_top - table_y - plaque_h - gap_before_plaque
    row_h = min(Inches(0.4), max_table_h // n_rows)
    row_h = max(row_h, Inches(0.26))
    font_scale = min(1.0, row_h / Inches(0.4))

    gframe = slide.shapes.add_table(n_rows, 5, table_x, table_y, table_w, row_h * n_rows)
    tbl = gframe.table
    for i, w in enumerate(col_w):
        tbl.columns[i].width = w

    def _cell(cell, text, size, color, bold, bg, align=PP_ALIGN.LEFT, strike=False):
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Pt(6)
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.name = FONT
        run.font.color.rgb = color if isinstance(color, RGBColor) else RGBColor.from_string(color)
        if strike:
            rPr = run._r.get_or_add_rPr()
            rPr.set("strike", "sngStrike")

    headers = ["№", "Ключевой системный элемент", "Обычная цена", "Скидка", "Цена по Программе лояльности"]
    aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER]
    h_size, h_size4 = 11.5 * font_scale, 10.5 * font_scale
    for c, (h, a) in enumerate(zip(headers, aligns)):
        _cell(tbl.cell(0, c), h, h_size if c != 4 else h_size4, WHITE, True, NAVY, a)

    body_size, final_size = 11 * font_scale, 12 * font_scale
    total_full, total_disc = 0, 0
    for i, item in enumerate(items):
        row = i + 1
        final_price = round(item["price"] * (1 - item["discount"] / 100))
        total_full += item["price"]; total_disc += final_price
        row_fill = WHITE if i % 2 == 0 else OFFWHITE
        _cell(tbl.cell(row, 0), str(i + 1), body_size, DARKTEXT, False, row_fill, PP_ALIGN.CENTER)
        _cell(tbl.cell(row, 1), item["name"], body_size, DARKTEXT, False, row_fill)
        _cell(tbl.cell(row, 2), _fmt_rub(item["price"]), body_size, TAUPE, False, row_fill, PP_ALIGN.CENTER,
              strike=item["discount"] > 0)
        disc_text = f'-{item["discount"]}%' if item["discount"] > 0 else "—"
        disc_bg = TERRACOTTA if item["discount"] > 0 else row_fill
        disc_color = WHITE if item["discount"] > 0 else DARKTEXT
        _cell(tbl.cell(row, 3), disc_text, body_size, disc_color, item["discount"] > 0, disc_bg, PP_ALIGN.CENTER)
        final_color = GOLD if item["discount"] > 0 else DARKTEXT
        _cell(tbl.cell(row, 4), _fmt_rub(final_price), final_size, final_color, True, row_fill, PP_ALIGN.CENTER)

    savings = total_full - total_disc
    plaque_y = table_y + row_h * n_rows + gap_before_plaque
    add_rounded_rect(slide, table_x, plaque_y, table_w, plaque_h, fill_color=RGBColor(0xFF, 0xF5, 0xE6),
                      line_color=GOLD, line_width_pt=1.5)
    add_textbox(slide, f"Ваша экономия при полном внедрении плана: {_fmt_rub(savings)}",
                table_x + Inches(0.3), plaque_y, table_w - Inches(0.6), plaque_h,
                size=13, color=DARKTEXT, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, client, dark=False)


def build_slide_bundle_fork(pres, data):
    client = data["client"]
    fork = data.get("bundleFork")
    if not fork:
        return  # на реальных данных развилка не всегда применима — слайд просто пропускается
    slide = add_slide(pres, NAVY)
    title(slide, "«Возрождение малого бизнеса»", dark=True)

    left_x, left_y, left_w, left_h = Inches(0.7), Inches(1.7), Inches(4.0), Inches(4.9)
    add_rounded_rect(slide, left_x, left_y, left_w, left_h, fill_color=RGBColor(0x0F, 0x3D, 0x51),
                      line_color=TEAL, line_width_pt=1)
    add_textbox(slide, "Системные элементы к внедрению в Вашем плане:", left_x + Inches(0.25), left_y + Inches(0.25),
                left_w - Inches(0.5), Inches(0.8), size=14, color=GOLD, bold=True)
    add_bullets(slide, fork["intersection"], left_x + Inches(0.25), left_y + Inches(1.1), left_w - Inches(0.4),
                left_h - Inches(1.3), size=13 if len(fork["intersection"]) <= 5 else 11, color=WHITE, line_spacing=1.3)

    right_x, right_w = Inches(5.6), Inches(5.0)
    gap = Inches(0.25)
    path_a_h, path_b_h = Inches(2.2), Inches(2.6)

    # Обе стрелки должны быть одинаковой длины и выходить из ОДНОЙ точки —
    # вертикального центра левой плашки «Системные элементы к внедрению».
    # Для этого центры плашек «Путь А» и «Путь Б» размещаются СИММЕТРИЧНО
    # относительно этой точки (на равном вертикальном расстоянии d от неё) —
    # тогда обе стрелки геометрически равны по длине без всяких допущений.
    origin_x, origin_y = left_x + left_w, left_y + left_h // 2
    d = gap // 2 + (path_a_h + path_b_h) // 4
    path_a_y = origin_y - d - path_a_h // 2
    path_b_y = origin_y + d - path_b_h // 2

    add_line_arrow(slide, origin_x, origin_y, right_x, path_a_y + path_a_h // 2, TAUPE, 4.5)
    add_line_arrow(slide, origin_x, origin_y, right_x, path_b_y + path_b_h // 2, TAUPE, 4.5)

    add_rounded_rect(slide, right_x, path_a_y, right_w, path_a_h, fill_color=RGBColor(0x12, 0x29, 0x3A),
                      line_color=TERRACOTTA, line_width_pt=1.25)
    add_textbox(slide, "ПУТЬ А — полные Программы по отдельности", right_x + Inches(0.25), path_a_y + Inches(0.14),
                right_w - Inches(0.5), Inches(0.5), size=13, color=TERRACOTTA, bold=True)
    add_bullets(slide, fork["intersection"], right_x + Inches(0.25), path_a_y + Inches(0.62), right_w - Inches(0.5),
                path_a_h - Inches(0.75), size=11, color=OFFWHITE, line_spacing=1.05)

    add_rounded_rect(slide, right_x, path_b_y, right_w, path_b_h, fill_color=RGBColor(0x12, 0x29, 0x3A),
                      line_color=GOLD, line_width_pt=1.25)
    add_textbox(slide, "ПУТЬ Б — Комплексная программа «Возрождение малого бизнеса»", right_x + Inches(0.25),
                path_b_y + Inches(0.15), right_w - Inches(0.5), Inches(0.65), size=13, color=GOLD, bold=True)
    add_bullets(slide, fork["bundlePrograms"], right_x + Inches(0.25), path_b_y + Inches(0.85), right_w - Inches(0.5),
                path_b_h - Inches(1.0), size=11, color=OFFWHITE, line_spacing=1.05)

    badge_x, badge_w = right_x + right_w + Inches(0.25), Inches(1.75)
    add_rounded_rect(slide, badge_x, path_a_y, badge_w, path_a_h, fill_color=RGBColor(0x1A, 0x1A, 0x1A),
                      line_color=TERRACOTTA, line_width_pt=1)
    add_textbox(slide, [f'{fork["individualWeeks"]} недель', _fmt_rub(fork["individualPrice"])],
                badge_x, path_a_y + path_a_h // 2 - Inches(0.5), badge_w, Inches(1.0),
                size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_rounded_rect(slide, badge_x, path_b_y, badge_w, path_b_h, fill_color=RGBColor(0x1A, 0x1A, 0x1A),
                      line_color=GOLD, line_width_pt=1)
    add_textbox(slide, [f'{fork["bundleWeeks"]} недель', _fmt_rub(fork["bundlePrice"])],
                badge_x, path_b_y + path_b_h // 2 - Inches(0.5), badge_w, Inches(1.0),
                size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    footer(slide, client, dark=True)


def build_slide_closing(pres, data):
    client = data["client"]
    slide = add_slide(pres, NAVY)
    add_textbox(slide, "ЧТО ДАЛЬШЕ?", Inches(1.5), Inches(1.5), W - Inches(3), Inches(0.7),
                size=30, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    steps = [
        "Выбираем Программы и порядок внедрения",
        "Фиксируем дату старта",
        "Договариваемся об организационных деталях",
    ]
    y = Inches(3.4)
    for i, step in enumerate(steps):
        add_rounded_rect(slide, W // 2 - Inches(3.6), y, Inches(0.5), Inches(0.5), fill_color=TERRACOTTA, radius=0.5)
        add_textbox(slide, str(i + 1), W // 2 - Inches(3.6), y, Inches(0.5), Inches(0.5),
                    size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, step, W // 2 - Inches(2.9), y, Inches(6.5), Inches(0.6),
                    size=16, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.85)
    footer(slide, client, dark=True)


def build_slide_mission(pres, data):
    slide = add_slide(pres, NAVY)
    logo_w = Inches(4.5)
    logo_h = int(logo_w * (221 / 357))
    from pptx_presentation_builder import LOGO_VERTICAL
    slide.shapes.add_picture(str(LOGO_VERTICAL), (W - logo_w) // 2, Inches(0.5), width=logo_w, height=logo_h)
    add_textbox(
        slide,
        ["Наша миссия – преумножать число", "устойчивых и растущих бизнесов"],
        Inches(1.5), Inches(0.5) + logo_h + Inches(0.5), W - Inches(3), Inches(0.9),
        size=22, color=WHITE, bold=True, italic=True, align=PP_ALIGN.CENTER,
    )
    add_textbox(slide, "www.fenix-lab.ru", Inches(1.5), Inches(0.5) + logo_h + Inches(1.9), W - Inches(3), Inches(0.4),
                size=16, color=TERRACOTTA, align=PP_ALIGN.CENTER)
