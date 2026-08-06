# -*- coding: utf-8 -*-
"""
pptx_presentation_builder.py — сборка Презентации для консультации на
python-pptx, ЦЕЛИКОМ в Python (замена pptx_build/*.js на pptxgenjs).

ПРИЧИНА ПЕРЕХОДА (06.08.2026): исходная реализация на pptxgenjs требовала
Node.js рядом с Python-бэкендом на Render — либо второй рантайм в том же
сервисе (Docker), либо отдельный сервис. Оба варианта добавляют инфраструктурную
сложность, которую Игорь сознательно предпочёл не брать на себя (минимум
поддержки, простота). Эта версия — третий модуль в том же семействе, что и
pdf_report_builder.py / consultation_script_builder.py: обычный Python-код,
без новых зависимостей рантайма, встраивается в app.py так же, как они.

СТАТУС: переносится ПОЭТАПНО. Список слайдов и их статус — см.
SLIDE_STATUS ниже. Слайды, которых здесь ещё нет, продолжают собираться
только через pptx_build/*.js (Node) — этот файл им пока не замена целиком.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = Path(__file__).parent
LOGO_VERTICAL = BASE / "assets" / "pptx_logos" / "fenix_logo_vertical.png"
LOGO_HORIZONTAL = BASE / "assets" / "pptx_logos" / "fenix_logo_horizontal.png"

# ---------------------------------------------------------------------------
# Фирменная палитра (та же, что в pptx_build/base.js и pdf_report_builder.py)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x30, 0x41)
TERRACOTTA = RGBColor(0xD5, 0x53, 0x0B)
GOLD = RGBColor(0xFC, 0xA7, 0x00)
TAUPE = RGBColor(0xB3, 0x92, 0x7E)
TEAL = RGBColor(0x1A, 0x75, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF7, 0xF5, 0xF2)
DARKTEXT = RGBColor(0x1A, 0x1A, 0x1A)

FONT = "Roboto"

W, H = Inches(13.33), Inches(7.5)  # LAYOUT_WIDE — тот же формат, что был на pptxgenjs

SLIDE_STATUS = {
    # slide_key: перенесён ли уже в Python (True) или пока только в Node (False)
    "cover": True,
    "agenda": False,
    "your_words": False,
    "reasons": False,
    "growth_rules": True,
    "kse_concept": False,
    "symptoms": False,
    "maturity_top5": True,
    "priority_chain": False,
    "cost_only": False,
    "full_contrast": False,
    "how_to_reach": False,
    "kse_list_repeat": False,
    "what_order": False,
    "house": False,
    "how_to_implement": False,
    "what_is_program": False,
    "program_model": False,
    "program_model_practice": False,
    "house_continuation": False,
    "programs_foundation": False,
    "programs_core": False,
    "programs_superstructure": False,
    "bundle_fork": False,
    "guarantee": False,
    "loyalty_program": False,
    "closing": False,
    "mission": False,
}


def new_presentation():
    pres = Presentation()
    pres.slide_width = W
    pres.slide_height = H
    return pres


def _blank_layout(pres):
    # Layout 6 в дефолтном шаблоне python-pptx — "Blank" (без плейсхолдеров)
    return pres.slide_layouts[6]


def add_slide(pres, bg_color=WHITE):
    slide = pres.slides.add_slide(_blank_layout(pres))
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    return slide


def _set_font(run, size, color, bold=False, italic=False, font_name=FONT):
    if isinstance(color, str):
        color = RGBColor.from_string(color)
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_textbox(slide, text, x, y, w, h, size, color, bold=False, italic=False,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name=FONT,
                 line_spacing=1.0, char_spacing=None):
    """text: строка ИЛИ список строк (каждая — отдельный параграф)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        _set_font(run, size, color, bold, italic, font_name)
        if char_spacing is not None:
            # межбуквенный интервал — через прямую правку XML (python-pptx не
            # даёт готового свойства), в сотых долях пункта
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(int(char_spacing * 100)))
    return box


def add_bullets(slide, items, x, y, w, h, size, color, font_name=FONT,
                 line_spacing=1.1, bold=False):
    """items: список строк — по одному буллету на строку (• в начале)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = f"•  {item}"
        _set_font(run, size, color, bold, False, font_name)
    return box


def add_rounded_rect(slide, x, y, w, h, fill_color, line_color=None, line_width_pt=1.0, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # radius (доля от меньшей стороны) — тот же смысл, что rectRadius у pptxgenjs
    try:
        shape.adjustments[0] = radius
    except IndexError:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    shape.shadow.inherit = False
    return shape


def add_line_arrow(slide, x1, y1, x2, y2, color, width_pt=1.5):
    """Стрелка из (x1,y1) в (x2,y2) — python-pptx поддерживает произвольные
    EMU-координаты у коннектора напрямую (в отличие от pptxgenjs, здесь не
    нужен трюк с min/abs/flip — connector умеет рисовать в любом направлении
    как есть)."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    # Наконечник стрелки на конце — прямая правка XML (python-pptx не даёт
    # готового свойства для头 arrowhead)
    ln = connector.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle"})
    ln.append(tail)
    return connector


def title(slide, text, dark=False):
    add_textbox(
        slide, text.upper(), Inches(0.7), Inches(0.5), W - Inches(1.4), Inches(0.9),
        size=30, color=WHITE if dark else NAVY, bold=True, font_name=FONT,
    )


def footer(slide, client, dark=False):
    slide.shapes.add_picture(str(LOGO_HORIZONTAL), Inches(0.5), H - Inches(0.53),
                              width=Inches(1.15), height=Inches(1.15 * (535 / 1919)))
    add_textbox(
        slide, "© ЛАБОРАТОРИЯ БИЗНЕС ЛИДЕРСТВА «ФЕНИКС»",
        Inches(0), H - Inches(0.45), W, Inches(0.3),
        size=9, color=RGBColor(0x9A, 0xAA, 0xB3) if dark else RGBColor(0x9A, 0x90, 0x88),
        align=PP_ALIGN.CENTER, char_spacing=0.5,
    )
    add_textbox(
        slide, f'{client["name"]} · {client["company"]}',
        W - Inches(6.5), H - Inches(0.45), Inches(6), Inches(0.3),
        size=9, color=RGBColor(0x9A, 0xAA, 0xB3) if dark else RGBColor(0x9A, 0x90, 0x88),
        align=PP_ALIGN.RIGHT,
    )
